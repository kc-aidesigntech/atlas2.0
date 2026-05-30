"""Lightweight PPTX slide previewer (no LibreOffice). Reads the actual shapes in a
generated deck and composites a PNG so we can visually verify layout/overlap.
Approximation: solid fills as rounded rects, first-run text, embedded pictures.
Usage: python3 _render_preview.py <deck.pptx> <slide_index> <out.png>"""
import sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

DPI = 96
def px(emu): return int(Emu(emu).inches * DPI)


def font(size_pt, bold=False):
    paths = [
        "/System/Library/Fonts/Helvetice.ttc",
        "/System/Library/Fonts/HelveticaNeue.ttc",
        "/Library/Fonts/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for pth in paths:
        try:
            return ImageFont.truetype(pth, int(size_pt * DPI / 72))
        except Exception:
            continue
    return ImageFont.load_default()


def rgb(c):
    return (c[0], c[1], c[2]) if c else (20, 20, 20)


def main():
    deck, idx, out = sys.argv[1], int(sys.argv[2]), sys.argv[3]
    prs = Presentation(deck)
    W, H = px(prs.slide_width), px(prs.slide_height)
    img = Image.new("RGB", (W, H), (0, 0, 0))
    d = ImageDraw.Draw(img)
    slide = list(prs.slides)[idx]
    for sh in slide.shapes:
        x, y, w, h = px(sh.left or 0), px(sh.top or 0), px(sh.width or 0), px(sh.height or 0)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic = Image.open(__import__("io").BytesIO(sh.image.blob)).convert("RGB")
                pic = pic.resize((max(1, w), max(1, h)))
                img.paste(pic, (x, y))
            except Exception as e:
                d.rectangle([x, y, x + w, y + h], outline=(120, 120, 120))
            continue
        try:
            if sh.fill.type is not None and sh.fill.fore_color and sh.fill.fore_color.type is not None:
                col = rgb(sh.fill.fore_color.rgb)
                if col != (0, 0, 0):
                    d.rounded_rectangle([x, y, x + w, y + h], radius=min(18, h // 2), fill=col)
        except Exception:
            pass
        if sh.has_text_frame and sh.text_frame.text.strip():
            ty = y
            for para in sh.text_frame.paragraphs:
                runs = para.runs
                if not runs:
                    continue
                tx = x
                size = runs[0].font.size.pt if runs[0].font.size else 14
                f = font(size, bool(runs[0].font.bold))
                text = "".join(r.text for r in runs)
                try:
                    col = rgb(runs[0].font.color.rgb)
                except Exception:
                    col = (255, 255, 255)
                if col == (0, 0, 0) and runs[0].font.color and False:
                    col = (255, 255, 255)
                d.text((tx, ty), text, fill=col if col != (0, 0, 0) else (10, 10, 10), font=f)
                ty += int(size * DPI / 72 * 1.25)
    img.save(out)
    print("wrote", out, img.size)


if __name__ == "__main__":
    main()
