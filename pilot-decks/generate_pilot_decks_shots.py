"""
Generate the second set of four Atlas pilot decks (.pptx) -- the "walkthrough"
edition -- where every step is paired with a real screenshot captured from the
live app (see pilot-decks/capture/capture_screens.py).

Reuses the theme/helpers from generate_pilot_decks.py (black surfaces, white text,
signal yellow + lucid green accents). Each step slide shows a compact
WHERE / DO / EXPECT instruction column on the left and the captured screen on the
right. Run: python3 pilot-decks/generate_pilot_decks_shots.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import generate_pilot_decks as base  # shared palette + slide helpers

ROOT = os.path.dirname(os.path.abspath(__file__))
SHOTS = os.path.join(ROOT, "screenshots")

# Pull theme constants from the base deck module so both sets stay identical.
B = base
YELLOW, GREEN, WHITE, GRAY, MUTED = B.YELLOW, B.GREEN, B.WHITE, B.GRAY, B.MUTED
BLACK, PANEL, PANEL_SOFT, HAIRLINE = B.BLACK, B.PANEL, B.PANEL_SOFT, B.HAIRLINE
FONT, EMU_W, EMU_H = B.FONT, B.EMU_W, B.EMU_H


def step_shot_slide(prs, accent, idx, total, title, where, do, expect, rpc, shot):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    B._set_bg(slide)
    B._rect(slide, 0, 0, Inches(0.22), EMU_H, accent)

    # Step number circle + overline + title (header band, full width).
    circ = B._rect(slide, Inches(0.6), Inches(0.55), Inches(0.78), Inches(0.78), accent, MSO_SHAPE.OVAL)
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(idx)
    r.font.name = FONT
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = BLACK
    _, tf = B._box(slide, Inches(1.6), Inches(0.5), Inches(11), Inches(0.34))
    B._runs(tf, f"STEP {idx} OF {total}", 11, accent, bold=True, first=True, spacing=2.4)
    _, tf = B._box(slide, Inches(1.6), Inches(0.82), Inches(11.4), Inches(0.7))
    B._runs(tf, title, 25, WHITE, bold=True, first=True, line_spacing=1.0)

    # Right: the live screenshot in a framed panel (16:9 source -> 7.66x4.31in).
    img_w = Inches(7.62)
    img_h = Inches(7.62 * 9 / 16)
    img_x = Inches(5.4)
    img_y = Inches(2.05)
    B._rect(slide, img_x - Inches(0.06), img_y - Inches(0.06), img_w + Inches(0.12), img_h + Inches(0.12),
            PANEL_SOFT, MSO_SHAPE.ROUNDED_RECTANGLE)
    path = os.path.join(SHOTS, shot)
    if os.path.exists(path):
        slide.shapes.add_picture(path, img_x, img_y, img_w, img_h)
    else:
        _, tf = B._box(slide, img_x, img_y + img_h / 2, img_w, Inches(0.5))
        B._runs(tf, f"[missing screenshot: {shot}]", 12, MUTED, first=True, align=PP_ALIGN.CENTER)
    _, tf = B._box(slide, img_x, img_y + img_h + Inches(0.12), img_w, Inches(0.32))
    B._runs(tf, "live capture \u00b7 localhost:5173", 9.5, MUTED, first=True, align=PP_ALIGN.RIGHT, spacing=0.6)

    # Left: compact WHERE / DO / EXPECT column (kept clear of the footer band).
    rows = [("WHERE", where, GREEN), ("DO", do, WHITE), ("EXPECT", expect, accent)]
    y = 2.1
    for label, text, lab_color in rows:
        _, tf = B._box(slide, Inches(0.6), Inches(y), Inches(4.5), Inches(0.3))
        B._runs(tf, label, 11, lab_color, bold=True, first=True, spacing=1.8)
        _, tf = B._box(slide, Inches(0.6), Inches(y + 0.26), Inches(4.55), Inches(1.0))
        B._runs(tf, text, 13.5, GRAY if label != "DO" else WHITE, first=True, line_spacing=1.1)
        y += 1.2
    if rpc:
        chip_y = 5.82
        B._rect(slide, Inches(0.6), Inches(chip_y), Inches(4.55), Inches(0.82), PANEL_SOFT,
                MSO_SHAPE.ROUNDED_RECTANGLE)
        _, tf = B._box(slide, Inches(0.8), Inches(chip_y + 0.13), Inches(4.2), Inches(0.3))
        B._runs(tf, "COMMAND RPC", 9, MUTED, bold=True, first=True, spacing=1.6)
        _, tf = B._box(slide, Inches(0.8), Inches(chip_y + 0.42), Inches(4.2), Inches(0.36))
        p, r = B._runs(tf, rpc, 12, GREEN, first=True)
        r.font.name = "Menlo"
    B._footer(slide, accent)
    return slide


CRED = "AtlasPilot2026!"
DECKS = {
    "Administrator": {
        "file": "Atlas-Pilot-Walkthrough-01-Administrator.pptx",
        "accent": YELLOW,
        "login": "pilot.admin@atlas.test",
        "subtitle": "Full-access governance. You see every enrollee and configure who can access what \u2014 while data-warehouse exports stay locked to the service role.",
        "steps": [
            ("Sign in as Administrator", "http://localhost:5173/app",
             "Enter the administrator login and password, then \u201csign in with email\u201d.",
             "The workspace bootstraps with the administrator shell.", None, "administrator/01_signin.png"),
            ("See every enrollee", "enrollees \u2192 enrollee picker",
             "Open the enrollee picker at the top of the workspace.",
             "All four enrollees are visible (Elena, Marcus, Sandra, test-reference1).",
             None, "administrator/02_enrollees.png"),
            ("Governance & access matrix", "admin controls \u2192 governance",
             "Open the system record control center to set roles and one-to-many assignments.",
             "Live counts load (4 enrollees, 12 people, 1 org); assignments are editable.",
             "fn_access_matrix_save_*", "administrator/03_governance.png"),
            ("Plan routes across partners", "care delivery \u2192 route planning",
             "Open an enrollee\u2019s route board to match Z-codes to partner stations.",
             "Ranked partner matches appear with score, match %, and station burden.",
             None, "administrator/05_route_planning.png"),
            ("Switch role view / identity", "account settings (top-right)",
             "Use the role switcher to preview any role and manage operator basics.",
             "Current role view = administrator; warehouse exports stay service-role only.",
             None, "administrator/06_account_settings.png"),
        ],
    },
    "Navigator": {
        "file": "Atlas-Pilot-Walkthrough-02-Navigator.pptx",
        "accent": GREEN,
        "login": "pilot.navigator@atlas.test",
        "subtitle": "Front-line case work, scoped to your assignments. You see only your two assigned enrollees and drive their journey end to end.",
        "steps": [
            ("Sign in as Navigator", "http://localhost:5173/app",
             "Enter the navigator login and password, then \u201csign in with email\u201d.",
             "The workspace opens scoped to your assignments.", None, "navigator/01_signin.png"),
            ("Only your assigned enrollees", "enrollees \u2192 enrollee picker",
             "Open the enrollee picker.",
             "Only Sandra Morrison and Marcus Thompson appear \u2014 your two assignments.",
             None, "navigator/02_enrollees.png"),
            ("Complete a burden survey", "open an enrollee \u2192 open burden survey",
             "Open the domain-spectrum burden survey for an enrollee and submit.",
             "The submission saves and history updates.",
             "fn_save_enrollee_burden_submission", "navigator/03_burden_survey.png"),
            ("Your profile & assignments", "my profile",
             "Review your assignment board and the enrollee burden surveys you own.",
             "Profile shows assigned enrollees: 2, with each claimed enrollment.",
             "fn_save_navigator_competency_assessment", "navigator/04_my_profile.png"),
            ("Refer someone", "refer",
             "Open the referral form and submit a referral into the intake queue.",
             "The referral enters the same in-system queue staff work from.",
             None, "navigator/05_refer.png"),
            ("Your station context", "my station",
             "Open your station to see the partner org and its Z-code coverage.",
             "Station context resolves for your linked partner organization.",
             "fn_get_my_navigator_station_context", "navigator/06_my_station.png"),
        ],
    },
    "Supervisor": {
        "file": "Atlas-Pilot-Walkthrough-03-Supervisor.pptx",
        "accent": GREEN,
        "login": "pilot.supervisor@atlas.test",
        "subtitle": "Oversight across the navigator team. You supervise every navigator and can see competency rollups and the enrollments your navigators cover.",
        "steps": [
            ("Sign in as Supervisor", "http://localhost:5173/app",
             "Enter the supervisor login and password, then \u201csign in with email\u201d.",
             "The workspace opens with the supervisor shell.", None, "supervisor/01_signin.png"),
            ("See your navigators", "assigned navigators",
             "Open the assigned-navigators board.",
             "The navigators that roll up to you are listed with their enrollee counts.",
             None, "supervisor/02_assigned_navigators.png"),
            ("Competency rollup", "navigator assessments",
             "Open the competency rollup for your navigators.",
             "Rolling weighted competency (3x / 2x / 1x by recency) is shown per navigator.",
             "fn_save_navigator_competency_assessment", "supervisor/03_navigator_assessments.png"),
            ("Team burden", "team burden",
             "Open the team burden view across your supervised navigators.",
             "Team-level burden metrics populate for your scope only.",
             None, "supervisor/04_team_burden.png"),
        ],
    },
    "Partner": {
        "file": "Atlas-Pilot-Walkthrough-04-Partner.pptx",
        "accent": YELLOW,
        "login": "pilot.partner@atlas.test",
        "subtitle": "Community partner self-service. You report service capacity for your organization \u2014 with zero access to enrollee health information.",
        "steps": [
            ("Sign in as Partner", "http://localhost:5173/app",
             "Enter the partner login and password, then \u201csign in with email\u201d.",
             "The workspace opens to your organization \u2014 no enrollee roster.",
             None, "partner/01_signin.png"),
            ("Your org capacity (no PHI)", "referral portal",
             "Review your organization\u2019s Z-code coverage and capacity overview.",
             "You see your org\u2019s Z-codes and radial scale \u2014 zero enrollee PHI.",
             None, "partner/02_capacity_overview.png"),
            ("Your station", "my station",
             "Open your station context for the organization.",
             "Station Z-code coverage and unresolved categories are shown.",
             None, "partner/03_my_station.png"),
            ("Service capacity survey", "service capacity",
             "Open the service-capacity survey history; start or edit a submission.",
             "Completed runs stay read-only; new submissions save for your org.",
             "fn_save_partner_service_capacity", "partner/04_service_capacity.png"),
            ("Public referral (logged out)", "http://localhost:5173/  (no login)",
             "Sign out and submit a referral from the public Referral Portal.",
             "It persists to the database for staff to claim later \u2014 the only anon write.",
             None, "partner/05_public_referral.png"),
        ],
    },
}


def build():
    for role, cfg in DECKS.items():
        prs = Presentation()
        prs.slide_width = EMU_W
        prs.slide_height = EMU_H
        accent = cfg["accent"]
        base.title_slide(prs, role, accent, cfg["subtitle"], cfg["login"], CRED)
        total = len(cfg["steps"])
        for i, (title, where, do, expect, rpc, shot) in enumerate(cfg["steps"], start=1):
            step_shot_slide(prs, accent, i, total, title, where, do, expect, rpc, shot)
        base.closing_slide(prs, accent, role)
        out = os.path.join(ROOT, cfg["file"])
        prs.save(out)
        n = len(prs.slides._sldIdLst)
        print(f"wrote {out} ({n} slides)")


if __name__ == "__main__":
    build()
