"""
Generate four themed Atlas pilot decks (.pptx), one per role, that lay out a
simple step-by-step pilot exploration plan correlated with PILOT.md.

Theme mirrors the app: black surfaces, white text, signal yellow (#FCC01A) and
lucid green (#81BC36) accents, "atlas x lucid living / Community Navigation
Platform" branding. Run: python3 pilot-decks/generate_pilot_decks.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette -------------------------------------------------------------
BLACK = RGBColor(0x00, 0x00, 0x00)
PANEL = RGBColor(0x0C, 0x0C, 0x0C)
PANEL_SOFT = RGBColor(0x14, 0x14, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xC7, 0xC7, 0xC7)
MUTED = RGBColor(0x8F, 0x8F, 0x8F)
YELLOW = RGBColor(0xFC, 0xC0, 0x1A)
GREEN = RGBColor(0x81, 0xBC, 0x36)
HAIRLINE = RGBColor(0x2A, 0x2A, 0x2A)

FONT = "Helvetica Neue"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _set_bg(slide, color=BLACK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def _runs(tf, text, size, color, bold=False, italic=False, spacing=None,
          first=False, align=PP_ALIGN.LEFT, line_spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if spacing is not None:
        _letter_spacing(r, spacing)
    return p, r


def _letter_spacing(run, pts):
    # Tracking is not exposed by python-pptx; set the spc attribute directly.
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    _no_line(sp)
    sp.shadow.inherit = False
    return sp


def _pill(slide, x, y, text, fill, txt_color):
    w = Inches(0.04 * len(text) + 0.5)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    _no_line(sp)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = txt_color
    _letter_spacing(r, 1.2)
    return sp


def _footer(slide, accent):
    _rect(slide, Inches(0.55), Inches(6.95), Inches(0.5), Pt(3), accent)
    _, tf = _box(slide, Inches(1.15), Inches(6.86), Inches(7), Inches(0.4))
    _runs(tf, "atlas \u00d7 lucid living   \u2014   Community Navigation Platform Pilot",
          10.5, MUTED, first=True, spacing=0.6)
    _, tf2 = _box(slide, Inches(10.5), Inches(6.86), Inches(2.3), Inches(0.4))
    _runs(tf2, "PILOT.md", 10.5, MUTED, first=True, align=PP_ALIGN.RIGHT)


def title_slide(prs, role, accent, subtitle, login, pw):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    # left accent column
    _rect(slide, 0, 0, Inches(0.22), EMU_H, accent)
    _, tf = _box(slide, Inches(0.9), Inches(1.5), Inches(11), Inches(0.5))
    _runs(tf, "ATLAS \u00d7 LUCID LIVING", 15, GREEN, bold=True, first=True, spacing=3.2)
    _, tf = _box(slide, Inches(0.86), Inches(2.15), Inches(11.6), Inches(2.2))
    _runs(tf, role, 54, WHITE, bold=True, first=True)
    _runs(tf, "Pilot Exploration Plan", 30, accent, bold=True)
    _, tf = _box(slide, Inches(0.9), Inches(4.5), Inches(11), Inches(0.7))
    _runs(tf, subtitle, 16, GRAY, first=True, line_spacing=1.2)
    # login card
    card = _rect(slide, Inches(0.9), Inches(5.35), Inches(7.6), Inches(1.15), PANEL,
                 MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = _box(slide, Inches(1.15), Inches(5.5), Inches(7.2), Inches(0.9))
    _runs(tf, "YOUR LOGIN", 10, accent, bold=True, first=True, spacing=2)
    _runs(tf, f"{login}      \u00b7      {pw}", 17, WHITE, bold=True)
    _footer(slide, accent)
    return slide


def step_slide(prs, accent, idx, total, title, where, do, expect, rpc=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _rect(slide, 0, 0, Inches(0.22), EMU_H, accent)
    # step number circle
    circ = _rect(slide, Inches(0.85), Inches(0.85), Inches(1.0), Inches(1.0), accent,
                 MSO_SHAPE.OVAL)
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(idx)
    r.font.name = FONT
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = BLACK
    # overline + title
    _, tf = _box(slide, Inches(2.1), Inches(0.92), Inches(10.5), Inches(0.4))
    _runs(tf, f"STEP {idx} OF {total}", 12, accent, bold=True, first=True, spacing=2.4)
    _, tf = _box(slide, Inches(2.1), Inches(1.28), Inches(10.6), Inches(1.1))
    _runs(tf, title, 30, WHITE, bold=True, first=True, line_spacing=1.0)
    # hairline
    _rect(slide, Inches(0.9), Inches(2.75), Inches(11.5), Pt(1), HAIRLINE)
    # labeled rows
    rows = [("WHERE", where, GREEN), ("DO", do, WHITE), ("EXPECT", expect, accent)]
    y = 3.05
    for label, text, lab_color in rows:
        _, tf = _box(slide, Inches(0.9), Inches(y), Inches(1.7), Inches(0.6))
        _runs(tf, label, 12, lab_color, bold=True, first=True, spacing=1.8)
        _, tf = _box(slide, Inches(2.7), Inches(y - 0.04), Inches(9.7), Inches(1.0))
        _runs(tf, text, 17, GRAY if label != "DO" else WHITE, first=True, line_spacing=1.15)
        y += 1.08
    if rpc:
        chip = _rect(slide, Inches(2.7), Inches(y + 0.02), Inches(9.4), Inches(0.62), PANEL_SOFT,
                     MSO_SHAPE.ROUNDED_RECTANGLE)
        _, tf = _box(slide, Inches(2.9), Inches(y + 0.13), Inches(9.0), Inches(0.4))
        p, _ = _runs(tf, "RPC  ", 11, MUTED, bold=True, first=True, spacing=1.5)
        r = p.add_run()
        r.text = rpc
        r.font.name = "Menlo"
        r.font.size = Pt(12)
        r.font.color.rgb = GREEN
    _footer(slide, accent)
    return slide


def closing_slide(prs, accent, role):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _rect(slide, 0, 0, Inches(0.22), EMU_H, accent)
    _, tf = _box(slide, Inches(0.9), Inches(1.4), Inches(11), Inches(0.4))
    _runs(tf, "WRAP-UP", 13, accent, bold=True, first=True, spacing=2.6)
    _, tf = _box(slide, Inches(0.86), Inches(1.85), Inches(11.6), Inches(1.0))
    _runs(tf, "Verify & reset", 40, WHITE, bold=True, first=True)
    notes = [
        ("Confirm alignment", "Run verification/pilot_verify.sql to replay every screen read under this identity and print an ok / row-count / error table \u2014 no clicking required."),
        ("Fail-loud by design", "A red banner means a real permission/RLS denial, not an empty screen. The workspace never hides a denial behind blank data."),
        ("Reset the pilot", "Re-run verification/pilot_setup.sql (idempotent) to recreate the four logins and their seeded scope at any time."),
    ]
    y = 3.1
    for head, body in notes:
        _rect(slide, Inches(0.9), Inches(y + 0.04), Inches(0.16), Inches(0.16), accent, MSO_SHAPE.OVAL)
        _, tf = _box(slide, Inches(1.25), Inches(y - 0.05), Inches(11), Inches(1.0))
        p, _ = _runs(tf, head + "  ", 16, WHITE, bold=True, first=True)
        r = p.add_run(); r.text = "\u2014 " + body
        r.font.name = FONT; r.font.size = Pt(15); r.font.color.rgb = GRAY
        y += 1.15
    _footer(slide, accent)
    return slide


# ---- per-role content ----------------------------------------------------
ROLES = {
    "Administrator": {
        "file": "Atlas-Pilot-01-Administrator.pptx",
        "accent": YELLOW,
        "login": "pilot.admin@atlas.test",
        "pw": "AtlasPilot2026!",
        "subtitle": "Full-access governance. You see every enrollee and configure who can access what \u2014 while data-warehouse exports stay locked to the service role.",
        "steps": [
            ("Sign in as Administrator", "http://localhost:5173/app, role switcher \u2192 Administrator",
             "Sign in, then open \u201cassigned enrollees\u201d.",
             "All 4 enrollees are visible (Sandra, Marcus, Elena, test-reference1).", None),
            ("Set a person\u2019s roles", "governance / system operations \u2192 Access Matrix",
             "Pick a person, change their role chips, Save.",
             "Roles persist and reload with the new assignment.",
             "fn_access_matrix_save_person_roles"),
            ("Assign a navigator to an enrollment", "Access Matrix \u2192 enrollment navigators",
             "Choose a navigator for an enrollment, Save.",
             "The enrollment shows the new navigator.",
             "fn_access_matrix_save_enrollment_navigators"),
            ("Assign supervisor \u2192 navigators", "Access Matrix \u2192 supervisors",
             "Link a supervisor to one or more navigators, Save.",
             "Supervision edges appear on the supervision board.",
             "fn_access_matrix_save_navigator_supervisors"),
            ("Set partner primary contacts", "Access Matrix \u2192 partner contacts",
             "Assign a primary contact to a partner org, Save.",
             "That contact can now access the partner scope.",
             "fn_access_matrix_save_partner_contacts"),
            ("Review Z-code survey history", "system operations",
             "Open the Z-code domain survey history (admin-only).",
             "Aggregated survey history loads.", None),
            ("Confirm the warehouse is hidden", "anywhere in the workspace",
             "Look for data-warehouse export views.",
             "None are reachable \u2014 they are service-role only.", None),
        ],
    },
    "Navigator": {
        "file": "Atlas-Pilot-02-Navigator.pptx",
        "accent": GREEN,
        "login": "pilot.navigator@atlas.test",
        "pw": "AtlasPilot2026!",
        "subtitle": "Front-line case work, scoped to your assignments. You see only your two assigned enrollees and drive their journey end to end.",
        "steps": [
            ("Sign in as Navigator", "http://localhost:5173/app, role switcher \u2192 Navigator",
             "Sign in and open the enrollee dropdown.",
             "Only Sandra Morrison and Marcus Thompson appear \u2014 your assignments.", None),
            ("Claim an enrollment", "requests to enroll",
             "Pick a queued enrollment and claim it to yourself.",
             "It moves into your assigned list.",
             "fn_navigator_assign_enrollment_to_self"),
            ("Complete an enrollee burden survey", "open an enrollee \u2192 burden survey",
             "Fill the domain-spectrum survey and Submit.",
             "Submission saves and history updates.",
             "fn_save_enrollee_burden_submission"),
            ("Take a regulation / renewal test", "Regulation Tests panel",
             "Complete the test and Submit.",
             "Result is recorded against your profile.",
             "fn_save_regulation_test_submission"),
            ("Record a competency self-assessment", "navigator assessments",
             "Score the competency rubric and save.",
             "Self-assessment is stored.",
             "fn_save_navigator_competency_assessment"),
            ("Intake inferred Z-codes", "enrollee intake",
             "Confirm the inferred Z-codes for an enrollment, save.",
             "Z-codes attach to the enrollment.",
             "fn_intake_enrollment_inferred_z_codes"),
            ("Resolve a Z-code to a partner", "enrollee Z-code strip \u2192 resolve",
             "Mark a Z-code resolved and attribute it to a partner.",
             "The Z-code shows resolved with partner attribution.",
             "fn_set_enrollee_z_code_resolution"),
        ],
    },
    "Supervisor": {
        "file": "Atlas-Pilot-03-Supervisor.pptx",
        "accent": GREEN,
        "login": "pilot.supervisor@atlas.test",
        "pw": "AtlasPilot2026!",
        "subtitle": "Oversight across the navigator team. You supervise every navigator and can see competency rollups and the enrollments your navigators cover.",
        "steps": [
            ("Sign in as Supervisor", "http://localhost:5173/app, role switcher \u2192 Supervisor",
             "Sign in and open \u201cassigned navigators\u201d.",
             "The navigators you supervise are listed.", None),
            ("Review the competency rollup", "navigator assessments / team burden",
             "Open the navigator competency rollup.",
             "Team competency and burden metrics populate.", None),
            ("Record a supervision assessment", "navigator assessments",
             "Capture a supervision session / competency review and save.",
             "The assessment is stored for that navigator.",
             "fn_save_navigator_competency_assessment"),
            ("Confirm scoped enrollee visibility", "assigned enrollees",
             "Open enrollees covered by your navigators.",
             "You see those enrollees; none outside your supervised navigators.", None),
        ],
    },
    "Partner": {
        "file": "Atlas-Pilot-04-Partner.pptx",
        "accent": YELLOW,
        "login": "pilot.partner@atlas.test",
        "pw": "AtlasPilot2026!",
        "subtitle": "Community partner self-service. You report service capacity for your organization \u2014 with zero access to enrollee health information.",
        "steps": [
            ("Sign in as Partner", "http://localhost:5173/app, role switcher \u2192 Partner",
             "Sign in and look for any enrollee roster.",
             "There is none \u2014 a partner sees zero enrollee PHI.", None),
            ("Submit the service capacity survey", "service capacity",
             "Complete the Z-code capacity survey and Submit.",
             "Submission saves for \u201cAtlas Pilot Partner Org\u201d.",
             "fn_save_partner_service_capacity"),
            ("Ensure a partner identifier", "service capacity flow",
             "Provide contact details during the capacity flow.",
             "A partner identifier record is ensured.",
             "fn_ensure_partner_identifier"),
            ("Nullify a survey answer", "capacity history",
             "Void a prior answer with a reason.",
             "The answer is marked nullified with an audit note.",
             "fn_set_partner_survey_answer_nullification"),
            ("Try a public referral (logged out)", "http://localhost:5173/  (no login)",
             "Open the Referral Portal and submit a referral.",
             "It persists to the database for staff to claim later.", None),
        ],
    },
}


def build():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    for role, cfg in ROLES.items():
        prs = Presentation()
        prs.slide_width = EMU_W
        prs.slide_height = EMU_H
        accent = cfg["accent"]
        title_slide(prs, role, accent, cfg["subtitle"], cfg["login"], cfg["pw"])
        total = len(cfg["steps"])
        for i, (t, where, do, expect, rpc) in enumerate(cfg["steps"], start=1):
            step_slide(prs, accent, i, total, t, where, do, expect, rpc)
        closing_slide(prs, accent, role)
        path = os.path.join(out_dir, cfg["file"])
        prs.save(path)
        print("wrote", path, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
